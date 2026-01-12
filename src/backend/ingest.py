"""
Document ingestion - extract text from multiple formats with improved error handling.
"""

import os
from pathlib import Path
from typing import List, Tuple, Optional
import re

from .logger_config import logger


class DocumentIngestor:
    """Handle document extraction and chunking with proper error handling."""
    
    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        enable_ocr: bool = False,
        chunking_level: Optional[int] = None
    ):
        """
        Initialize document ingestor.
        
        Args:
            chunk_size: Size of text chunks in characters
            chunk_overlap: Overlap between chunks in characters
            enable_ocr: Enable OCR for image files (requires pytesseract)
        """
        if chunk_size <= 0:
            raise ValueError("chunk_size must be positive")
        if chunk_overlap < 0:
            raise ValueError("chunk_overlap must be non-negative")
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be less than chunk_size")
        
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.enable_ocr = enable_ocr
        self.chunking_level: Optional[int] = None
        # Store stats about the last chunking operation so the API/UI can report patterns used
        self.last_chunk_stats: dict = {}
        # Special terms that should be kept together
        self.preserve_terms = ['m2', 'M2', 'mileage', 'allowance', 'transportation', 'benefits']
        logger.info(f"DocumentIngestor initialized: chunk_size={chunk_size}, overlap={chunk_overlap}, ocr={enable_ocr}")
        if chunking_level is not None:
            self.set_chunking_level(chunking_level)

    def set_chunking_level(self, level: int) -> None:
        """Update chunk size and overlap according to a 1-10 level scale."""
        level = max(1, min(10, int(level)))

        # Map slider level to practical chunk parameters
        min_chunk = 350
        max_chunk = 2000
        min_overlap = 40
        max_overlap = 400

        ratio = 0.0 if level == 1 else (level - 1) / 9
        computed_chunk = int(min_chunk + ratio * (max_chunk - min_chunk))
        computed_overlap = int(min_overlap + ratio * (max_overlap - min_overlap))

        # Ensure chunk overlap never exceeds chunk size minus a buffer
        max_allowed_overlap = max(0, computed_chunk - 100)
        computed_overlap = min(computed_overlap, max_allowed_overlap)

        if computed_overlap <= 0:
            computed_overlap = min_overlap

        self.chunking_level = level
        self.chunk_size = max(min_chunk, computed_chunk)
        self.chunk_overlap = max(0, min(computed_overlap, self.chunk_size - 1))

        logger.info(
            "Chunking level set to %s (chunk_size=%s, chunk_overlap=%s)",
            self.chunking_level,
            self.chunk_size,
            self.chunk_overlap
        )
    
    def load_and_process_documents(self, file_paths: List[str]) -> Tuple[List[str], str]:
        """
        Load and chunk documents from file paths with comprehensive logging.
        
        Args:
            file_paths: List of file paths to process
            
        Returns:
            Tuple of (chunks, combined_document_name)
        """
        logger.info(f"=== Starting document processing flow for {len(file_paths)} file(s) ===")
        
        # Step 1: Validate input
        if not file_paths:
            logger.warning("STEP 1 FAILED: No file paths provided")
            return [], "unknown"
        logger.info(f"STEP 1 COMPLETE: Validated {len(file_paths)} file path(s)")
        
        # Step 2: Process each file
        all_text = ""
        doc_names = []
        processed_count = 0
        failed_count = 0
        
        for idx, file_path in enumerate(file_paths, 1):
            logger.info(f"STEP 2.{idx}: Processing file {idx}/{len(file_paths)}: {file_path}")
            
            # Check file existence
            if not os.path.exists(file_path):
                logger.warning(f"STEP 2.{idx} FAILED: File not found: {file_path}")
                failed_count += 1
                continue
            logger.debug(f"STEP 2.{idx}.1 COMPLETE: File exists")
            
            # Extract text
            try:
                logger.debug(f"STEP 2.{idx}.2: Extracting text from {Path(file_path).suffix}")
                text = self._extract_text(file_path)
                
                if text and text.strip():
                    all_text += "\n\n" + text
                    doc_names.append(Path(file_path).stem)
                    processed_count += 1
                    logger.info(f"STEP 2.{idx} COMPLETE: Extracted {len(text)} characters from {file_path}")
                else:
                    logger.warning(f"STEP 2.{idx} FAILED: No text extracted from {file_path}")
                    failed_count += 1
            except Exception as e:
                logger.error(f"STEP 2.{idx} FAILED: Error processing {file_path}: {e}", exc_info=True)
                failed_count += 1
                continue
        
        logger.info(f"STEP 2 COMPLETE: Processed {processed_count} file(s), {failed_count} failed")
        
        # Step 3: Validate extracted text
        if not all_text.strip():
            logger.warning("STEP 3 FAILED: No text extracted from any files")
            return [], "unknown"
        logger.info(f"STEP 3 COMPLETE: Total extracted text length: {len(all_text)} characters")
        
        # Step 4: Chunk text
        logger.info("STEP 4: Starting text chunking process")
        try:
            chunks = self._chunk_text(all_text)
            if not chunks:
                logger.warning("STEP 4 FAILED: No chunks created from text")
                return [], "unknown"
            logger.info(f"STEP 4 COMPLETE: Created {len(chunks)} chunk(s)")
        except Exception as e:
            logger.error(f"STEP 4 FAILED: Error during chunking: {e}", exc_info=True)
            return [], "unknown"
        
        # Step 5: Prepare result
        doc_name = " + ".join(doc_names) if doc_names else "unknown"
        logger.info(f"=== Document processing flow COMPLETE: {len(chunks)} chunks from {processed_count} document(s) ===")
        
        return chunks, doc_name
    
    def _extract_text(self, file_path: str) -> str:
        """
        Extract text based on file type with logging.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Extracted text
        """
        ext = Path(file_path).suffix.lower()
        logger.debug(f"Determining extractor for file type: {ext}")
        
        extractors = {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_docx,
            '.txt': self._extract_text_file,
            '.md': self._extract_text_file,
            '.csv': self._extract_csv,
            '.xlsx': self._extract_excel,
            '.xls': self._extract_excel_legacy,
            '.pptx': self._extract_pptx,
            '.html': self._extract_html,
            '.htm': self._extract_html,
            '.xml': self._extract_xml,
            '.png': self._extract_image,
            '.jpg': self._extract_image,
            '.jpeg': self._extract_image,
        }
        
        extractor = extractors.get(ext)
        if extractor:
            logger.debug(f"Using extractor for {ext} format")
            result = extractor(file_path)
            if result:
                logger.debug(f"Extraction successful for {ext}: {len(result)} characters")
            else:
                logger.warning(f"Extraction returned empty result for {ext}")
            return result
        else:
            logger.warning(f"EXTRACTION FAILED: Unsupported file type: {ext}")
            return ""
    
    def _extract_pdf(self, file_path: str) -> str:
        """
        Extract text from PDF file using Docling for better layout handling,
        with fallback to PyPDF2 for simple cases.
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        # Try Docling first for complex layouts (infographics, tables, multi-column)
        try:
            from docling.document_converter import DocumentConverter
            
            logger.debug(f"Attempting Docling extraction for PDF: {file_path}")
            converter = DocumentConverter()
            result = converter.convert(file_path)
            
            # Extract text with structure preservation
            text = result.document.export_to_markdown()
            
            if text and len(text.strip()) > 50:  # Reasonable content threshold
                logger.info(f"Successfully extracted PDF with Docling: {file_path}")
                return text.strip()
            else:
                logger.warning(f"Docling returned minimal content, trying fallback for {file_path}")
                
        except Exception as e:
            logger.warning(f"Docling extraction failed, using PyPDF2 fallback: {e}")
        
        # Fallback to PyPDF2 for simple PDFs or if Docling fails
        try:
            from PyPDF2 import PdfReader
            
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PdfReader(f)
                num_pages = len(pdf_reader.pages)
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception as e:
                        logger.warning(f"Error extracting page {page_num} from {file_path}: {e}")
                        continue
                
                logger.debug(f"Extracted {num_pages} pages from PDF using PyPDF2: {file_path}")
            
            return text.strip()
        except ImportError:
            logger.error("PyPDF2 not installed. Install with: pip install PyPDF2")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PDF {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_pdf_with_ocr(self, file_path: str) -> str:
        """Extract text from graphics-heavy PDFs by converting to images and using OCR."""
        try:
            from pdf2image import convert_from_path
            import pytesseract
            
            logger.info(f"Converting PDF to images for OCR: {file_path}")
            images = convert_from_path(file_path, dpi=300)
            
            text = ""
            for i, image in enumerate(images, 1):
                page_text = pytesseract.image_to_string(image, lang='eng')
                if page_text.strip():
                    text += f"\n--- Page {i} ---\n{page_text}\n"
            
            logger.info(f"OCR extracted {len(text)} chars from {len(images)} pages")
            return text.strip()
            
        except ImportError:
            logger.error("pdf2image or pytesseract not installed")
            logger.error("Install: pip install pdf2image pytesseract")
            logger.error("Also install Tesseract: https://github.com/UB-Mannheim/tesseract/wiki")
            return ""
        except Exception as e:
            logger.error(f"PDF OCR failed: {e}", exc_info=True)
            return ""
    
    def _extract_docx(self, file_path: str) -> str:
        """
        Extract text from DOCX file.
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Extracted text
        """
        try:
            from docx import Document
            
            text = ""
            doc = Document(file_path)
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            logger.debug(f"Extracted text from DOCX: {file_path}")
            return text.strip()
        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return ""
        except Exception as e:
            logger.error(f"Error extracting DOCX {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_text_file(self, file_path: str) -> str:
        """
        Extract text from TXT or Markdown file.
        
        Args:
            file_path: Path to text file
            
        Returns:
            File contents
        """
        try:
            # Try UTF-8 first, fallback to other encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                        logger.debug(f"Extracted text from {file_path} using {encoding}")
                        return content
                except UnicodeDecodeError:
                    continue
            
            logger.error(f"Could not decode {file_path} with any encoding")
            return ""
        except Exception as e:
            logger.error(f"Error reading text file {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_csv(self, file_path: str) -> str:
        """Extract text from CSV file with hybrid row/column format."""
        try:
            import pandas as pd
            
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            df = None
            
            for encoding in encodings:
                try:
                    df = pd.read_csv(file_path, encoding=encoding, on_bad_lines='skip')
                    break
                except (UnicodeDecodeError, pd.errors.ParserError):
                    continue
            
            if df is None:
                logger.error(f"Could not read CSV {file_path}")
                return ""
            
            text = f"CSV File: {Path(file_path).name}\n\n"
            
            # Include column letters with headers (compact format)
            columns = df.columns.astype(str).tolist()
            column_letters = [chr(65 + i) if i < 26 else f"Col{i+1}" for i in range(len(columns))]
            text += "Columns: " + " | ".join([f"{letter}:{name}" for letter, name in zip(column_letters, columns)]) + "\n\n"
            
            # Extract rows with compact row references [R#]
            for idx, row in df.iterrows():
                row_num = idx + 2  # +2 because CSV row 1 is headers, pandas is 0-indexed
                row_data = []
                for val in row:
                    if pd.notna(val) and str(val).strip():
                        row_data.append(str(val))
                
                if row_data:
                    text += f"[R{row_num}] " + " | ".join(row_data) + "\n"
                
                # Repeat column headers every 25 rows for chunk context
                if (idx + 1) % 25 == 0 and idx > 0:
                    text += "\n[Columns: " + " | ".join(columns) + "]\n\n"
            
            logger.debug(f"Extracted {len(df)} rows from CSV with hybrid format: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("pandas not installed. Install with: pip install pandas")
            return ""
        except Exception as e:
            logger.error(f"Error extracting CSV {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_excel(self, file_path: str) -> str:
        """Extract text from Excel (.xlsx) file with hybrid row/column format."""
        try:
            import openpyxl
            from openpyxl.utils import get_column_letter
            
            text = f"Excel File: {Path(file_path).name}\n\n"
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n=== Sheet: {sheet_name} ===\n\n"
                
                # Extract headers with column letters (compact format)
                if sheet.max_row > 0:
                    headers = []
                    for col_idx, cell in enumerate(sheet[1], start=1):
                        col_letter = get_column_letter(col_idx)
                        header = str(cell.value) if cell.value else f"Col_{col_letter}"
                        headers.append((col_letter, header))
                    
                    if headers:
                        text += "Columns: " + " | ".join([f"{letter}:{name}" for letter, name in headers]) + "\n\n"
                
                # Extract data rows with compact row references [R#]
                for row_idx, row in enumerate(sheet.iter_rows(min_row=2, values_only=True), start=2):
                    row_data = []
                    for i, val in enumerate(row):
                        if i < len(headers) and val is not None and str(val).strip():
                            row_data.append(str(val))
                    
                    if row_data:
                        text += f"[R{row_idx}] " + " | ".join(row_data) + "\n"
                    
                    # Repeat column headers every 25 rows for chunk context
                    if (row_idx - 1) % 25 == 0 and row_idx > 2:
                        text += "\n[Columns: " + " | ".join([h[1] for h in headers]) + "]\n\n"
            
            logger.debug(f"Extracted Excel file with {len(workbook.sheetnames)} sheets with hybrid format: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("openpyxl not installed. Install with: pip install openpyxl")
            return ""
        except Exception as e:
            logger.error(f"Error extracting Excel {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_excel_legacy(self, file_path: str) -> str:
        """Extract text from legacy Excel (.xls) file."""
        try:
            import xlrd
            
            text = f"Excel File (Legacy): {Path(file_path).name}\n\n"
            workbook = xlrd.open_workbook(file_path)
            
            for sheet_idx in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_idx)
                text += f"\n=== Sheet: {sheet.name} ===\n\n"
                
                if sheet.nrows > 0:
                    headers = [str(sheet.cell_value(0, col)) for col in range(sheet.ncols)]
                    text += "Headers: " + " | ".join(headers) + "\n\n"
                
                for row_idx in range(1, sheet.nrows):
                    row = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
                    row_text = " | ".join(row)
                    if row_text.strip():
                        text += row_text + "\n"
            
            logger.debug(f"Extracted legacy Excel file: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("xlrd not installed. Install with: pip install xlrd")
            return ""
        except Exception as e:
            logger.error(f"Error extracting legacy Excel {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint (.pptx) file."""
        try:
            from pptx import Presentation
            
            text = f"PowerPoint File: {Path(file_path).name}\n\n"
            prs = Presentation(file_path)
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                text += f"\n--- Slide {slide_idx} ---\n\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
                    
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            if row_text.strip():
                                text += row_text + "\n"
            
            logger.debug(f"Extracted {len(prs.slides)} slides from PPTX: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return ""
        except Exception as e:
            logger.error(f"Error extracting PPTX {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_html(self, file_path: str) -> str:
        """Extract text from HTML file."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text(separator='\n', strip=True)
            text = re.sub(r'\n\s*\n', '\n\n', text)
            
            logger.debug(f"Extracted text from HTML: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("beautifulsoup4 not installed. Install with: pip install beautifulsoup4 lxml")
            return ""
        except Exception as e:
            logger.error(f"Error extracting HTML {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_xml(self, file_path: str) -> str:
        """Extract text from XML file."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                xml_content = f.read()
            
            soup = BeautifulSoup(xml_content, 'xml')
            text = soup.get_text(separator='\n', strip=True)
            text = re.sub(r'\n\s*\n', '\n\n', text)
            
            logger.debug(f"Extracted text from XML: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("beautifulsoup4 not installed. Install with: pip install beautifulsoup4 lxml")
            return ""
        except Exception as e:
            logger.error(f"Error extracting XML {file_path}: {e}", exc_info=True)
            return ""
    
    def _extract_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        if not self.enable_ocr:
            logger.warning(f"OCR disabled. Skipping image: {file_path}")
            return ""
        
        try:
            import pytesseract
            from PIL import Image
            
            # Configure Tesseract path if not in PATH
            # Uncomment and adjust if Tesseract is not in your system PATH:
            # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
            
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            
            logger.debug(f"Extracted text from image using OCR: {file_path}")
            return text.strip()
            
        except ImportError:
            logger.error("pytesseract/Pillow not installed. Install with: pip install pytesseract Pillow")
            logger.error("Also install Tesseract OCR: https://github.com/tesseract-ocr/tesseract")
            return ""
        except Exception as e:
            logger.error(f"Error extracting image {file_path}: {e}", exc_info=True)
            return ""
    
    def _chunk_text(self, text: str) -> List[str]:
        """
        Pattern-aware chunking:
        - Detects table-like, key-value/log-like, code, heading, list, and paragraph blocks using line heuristics.
        - Applies different chunking strategies per block type.
        - Stores stats about detected patterns in self.last_chunk_stats.
        """
        logger.debug("Starting pattern-aware text chunking process")

        # Reset stats for this run
        stats = {
            "total_blocks": 0,
            "paragraph_blocks": 0,
            "table_blocks": 0,
            "kv_blocks": 0,
            "code_blocks": 0,
            "heading_blocks": 0,
            "list_blocks": 0,
            "total_chunks": 0,
            "paragraph_chunks": 0,
            "table_chunks": 0,
            "kv_chunks": 0,
            "code_chunks": 0,
            "heading_chunks": 0,
            "list_chunks": 0,
            "chunking_level": self.chunking_level,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
        }

        # Step 1: Validate input
        if not text or not text.strip():
            logger.warning("CHUNKING STEP 1 FAILED: Empty text provided")
            self.last_chunk_stats = stats
            return []
        logger.debug(f"CHUNKING STEP 1 COMPLETE: Text validated ({len(text)} chars)")

        # --- Helper: classify a single line ---
        kv_patterns = [
            re.compile(r'^[^:]{1,60}:\s+.+$'),          # "Key: Value"
            re.compile(r'^\w[\w\s\-]{0,40}\s*=\s*.+$'), # "key = value"
        ]

        def classify_line(line: str) -> str:
            """Return 'blank', 'table', 'kv', 'code', 'heading', 'list', or 'normal'."""
            stripped = line.strip()
            if not stripped:
                return "blank"

            # Table row detection FIRST (before kv detection)
            # Check for Excel/CSV row format: [R#] or [Columns: ...]
            if re.match(r'^\[R\d+\]', stripped) or stripped.startswith('[Columns:'):
                return "table"

            # Heading detection: markdown-style or uppercase titles
            if stripped.startswith("#") or stripped.startswith("==") or stripped.startswith("--"):
                return "heading"
            if len(stripped) < 100 and stripped.isupper() and len(stripped.split()) <= 6:
                return "heading"

            # List detection: bullet points, numbered lists
            if stripped[0] in "-•*" or re.match(r'^\d+[\.\)]\s+', stripped) or stripped.startswith("- ") or stripped.startswith("* "):
                return "list"

            # Code detection: indentation, code markers, or language identifiers
            if stripped.startswith(("def ", "class ", "function ", "import ", "from ", "if __name__")):
                return "code"
            if stripped.startswith(("```", "~~~")) or stripped.startswith("\t") or (len(stripped) > 0 and stripped[0] == " " * 4):
                return "code"
            if re.match(r'^(import|from|def|class|async|await|return|try|except|finally)\b', stripped):
                return "code"

            # Table-ish: many separators or tabs
            if "|" in stripped or ("\t" in stripped):
                return "table"
            comma_count = stripped.count(",")
            if comma_count >= 3 and len(stripped) / (comma_count + 1) < 40:
                return "table"

            # Key–value / log-ish
            for pat in kv_patterns:
                if pat.match(stripped):
                    return "kv"

            return "normal"

        # --- Helper: group lines into typed blocks ---
        blocks = []  # each: {"type": "table"/"kv"/"paragraph", "lines": [...]}
        current_type: Optional[str] = None
        current_lines: List[str] = []

        lines = text.splitlines()

        def flush_block():
            nonlocal current_type, current_lines
            if current_type and current_lines:
                blocks.append({"type": current_type, "lines": current_lines})
            current_type = None
            current_lines = []

        for raw_line in lines:
            line_type = classify_line(raw_line)

            if line_type == "blank":
                # Blank line ends paragraph/kv blocks
                flush_block()
                continue

            # Normalize table/kv/paragraph types
            block_type = line_type
            if block_type == "normal":
                block_type = "paragraph"

            if current_type is None:
                current_type = block_type
                current_lines = [raw_line]
            elif block_type == current_type:
                current_lines.append(raw_line)
            else:
                flush_block()
                current_type = block_type
                current_lines = [raw_line]

        flush_block()  # flush last

        stats["total_blocks"] = len(blocks)
        for b in blocks:
            if b["type"] == "paragraph":
                stats["paragraph_blocks"] += 1
            elif b["type"] == "table":
                stats["table_blocks"] += 1
            elif b["type"] == "kv":
                stats["kv_blocks"] += 1
            elif b["type"] == "code":
                stats["code_blocks"] += 1
            elif b["type"] == "heading":
                stats["heading_blocks"] += 1
            elif b["type"] == "list":
                stats["list_blocks"] += 1

        logger.debug(f"CHUNKING STEP 2 COMPLETE: Grouped text into {len(blocks)} block(s)")

        # --- Helper: generic long-text chunker (sentence/word-aware, based on previous logic) ---
        def chunk_long_text(t: str) -> List[str]:
            t = re.sub(r'\s+', ' ', t).strip()
            if not t:
                return []
            if len(t) <= self.chunk_size:
                return [t]

            chunks_local: List[str] = []
            start = 0
            chunk_count_local = 0
            max_iterations_local = len(t) // (self.chunk_size - self.chunk_overlap) + 10

            while start < len(t) and chunk_count_local < max_iterations_local:
                chunk_count_local += 1
                end = start + self.chunk_size

                if end < len(t):
                    # Special handling for preserve terms - extend window if needed
                    for term in self.preserve_terms:
                        term_pos = t.find(term, start, end + 100)  # Look ahead slightly
                        if term_pos != -1 and term_pos > end - 50:  # If term is near boundary
                            # Find end of the term's context (next sentence or paragraph)
                            context_end = t.find('.', term_pos)
                            if context_end != -1 and context_end < end + 200:
                                end = context_end + 1
                                break
                    
                    # Normal sentence boundary detection
                    last_period = t.rfind('.', start, end)
                    last_exclamation = t.rfind('!', start, end)
                    last_question = t.rfind('?', start, end)
                    sentence_end = max(last_period, last_exclamation, last_question)

                    if sentence_end > start:
                        end = sentence_end + 1
                    else:
                        last_space = t.rfind(' ', start, end)
                        if last_space > start:
                            end = last_space

                chunk_local = t[start:end].strip()
                if chunk_local:
                    chunks_local.append(chunk_local)

                start = max(start + 1, end - self.chunk_overlap)
                if start >= len(t):
                    break

            return chunks_local

        # --- Helpers: per-block-type chunking ---
        def chunk_paragraph_block(block_lines: List[str]) -> List[str]:
            block_text = " ".join(l.strip() for l in block_lines if l.strip())
            return chunk_long_text(block_text)

        def chunk_table_block(block_lines: List[str]) -> List[str]:
            """Keep table structure; split by row count + char length with header preservation."""
            # Adaptive max_rows based on line verbosity
            if block_lines:
                avg_line_length = sum(len(l) for l in block_lines) / len(block_lines)
                if avg_line_length > 100:
                    max_rows = 10  # Verbose format
                elif avg_line_length > 60:
                    max_rows = 15
                else:
                    max_rows = 25  # Compact format
            else:
                max_rows = 25
            
            chunks_local: List[str] = []
            current_rows: List[str] = []
            
            # Extract header line (Columns: ...) if present
            header_line = None
            data_lines = block_lines
            if block_lines and block_lines[0].strip().startswith('Columns:'):
                header_line = block_lines[0]
                data_lines = block_lines[1:]

            def flush_rows():
                nonlocal current_rows
                if current_rows:
                    # Prepend header to each chunk for context
                    chunk_lines = [header_line] + current_rows if header_line else current_rows
                    joined = "\n".join(chunk_lines).strip()
                    if joined:
                        # if very large, fall back to generic chunker
                        if len(joined) > self.chunk_size * 2:
                            chunks_local.extend(chunk_long_text(joined))
                        else:
                            chunks_local.append(joined)
                current_rows = []

            for row in data_lines:
                current_rows.append(row)
                if len(current_rows) >= max_rows:
                    flush_rows()

            flush_rows()
            return chunks_local

        def chunk_kv_block(block_lines: List[str]) -> List[str]:
            """Group small key–value / log lines into compact chunks."""
            max_chars = max(self.chunk_size // 2, 400)
            chunks_local: List[str] = []
            current: List[str] = []
            current_len = 0

            for l in block_lines:
                s = l.strip()
                if not s:
                    continue
                if current_len + len(s) + 1 > max_chars:
                    if current:
                        chunks_local.append("\n".join(current).strip())
                    current = [s]
                    current_len = len(s)
                else:
                    current.append(s)
                    current_len += len(s) + 1

            if current:
                chunks_local.append("\n".join(current).strip())

            # For any chunk still too large, fall back to generic chunker
            final_chunks: List[str] = []
            for c in chunks_local:
                if len(c) > self.chunk_size:
                    final_chunks.extend(chunk_long_text(c))
                else:
                    final_chunks.append(c)
            return final_chunks

        # --- Apply per-block strategy ---
        all_chunks: List[str] = []
        for i, block in enumerate(blocks, 1):
            b_type = block["type"]
            b_lines = block["lines"]

            if b_type == "table":
                logger.debug(f"CHUNKING: Processing table block #{i} with {len(b_lines)} line(s)")
                table_chunks = chunk_table_block(b_lines)
                stats["table_chunks"] += len(table_chunks)
                all_chunks.extend(table_chunks)
            elif b_type == "kv":
                logger.debug(f"CHUNKING: Processing key-value block #{i} with {len(b_lines)} line(s)")
                kv_chunks = chunk_kv_block(b_lines)
                stats["kv_chunks"] += len(kv_chunks)
                all_chunks.extend(kv_chunks)
            else:  # paragraph
                logger.debug(f"CHUNKING: Processing paragraph block #{i} with {len(b_lines)} line(s)")
                para_chunks = chunk_paragraph_block(b_lines)
                stats["paragraph_chunks"] += len(para_chunks)
                all_chunks.extend(para_chunks)

        stats["total_chunks"] = len(all_chunks)

        # Derived helper fields for UI
        patterns = []
        if stats["paragraph_blocks"]:
            patterns.append("paragraph")
        if stats["table_blocks"]:
            patterns.append("table")
        if stats["kv_blocks"]:
            patterns.append("kv")
        if stats["code_blocks"]:
            patterns.append("code")
        if stats["heading_blocks"]:
            patterns.append("heading")
        if stats["list_blocks"]:
            patterns.append("list")
        stats["patterns"] = patterns

        self.last_chunk_stats = stats

        logger.info(
            f"CHUNKING COMPLETE: Created {len(all_chunks)} chunk(s) from "
            f"{len(blocks)} block(s) and {len(text)} characters"
        )
        return all_chunks
    
    def process_uploaded_file(self, file_content: bytes, filename: str) -> Tuple[List[str], str, dict]:
        """
        Process uploaded file from bytes with comprehensive logging.
        
        Args:
            file_content: File content as bytes
            filename: Original filename
            
        Returns:
            Tuple of (chunks, document_name, chunk_stats)
        """
        logger.info(f"=== Starting upload processing flow for file: {filename} ===")
        
        # Step 1: Validate file content
        if not file_content:
            logger.warning("STEP 1 FAILED: Empty file content provided")
            return [], "unknown", {}
        logger.info(f"STEP 1 COMPLETE: File content validated ({len(file_content)} bytes)")
        
        # Step 2: Create temp directory
        temp_path = Path("data/documents") / filename
        try:
            temp_path.parent.mkdir(parents=True, exist_ok=True)
            logger.info(f"STEP 2 COMPLETE: Temporary directory ready: {temp_path.parent}")
        except Exception as e:
            logger.error(f"STEP 2 FAILED: Cannot create temp directory: {e}", exc_info=True)
            return [], filename or "unknown", {}
        
        # Step 3: Write temporary file
        try:
            logger.info(f"STEP 3: Writing temporary file: {temp_path}")
            with open(temp_path, 'wb') as f:
                f.write(file_content)
            logger.info(f"STEP 3 COMPLETE: Temporary file written successfully")
        except Exception as e:
            logger.error(f"STEP 3 FAILED: Error writing temp file: {e}", exc_info=True)
            return [], filename or "unknown", {}
        
        # Step 4: Process document
        try:
            logger.info(f"STEP 4: Processing uploaded file: {filename}")
            chunks, doc_name = self.load_and_process_documents([str(temp_path)])
            stats = getattr(self, "last_chunk_stats", {}) or {}

            if chunks:
                logger.info(f"STEP 4 COMPLETE: Successfully processed into {len(chunks)} chunk(s)")
            else:
                logger.warning("STEP 4 FAILED: No chunks generated from file")

            return chunks, doc_name, stats
            
        except Exception as e:
            logger.error(f"STEP 4 FAILED: Error processing uploaded file {filename}: {e}", exc_info=True)
            return [], filename or "unknown", {}
        finally:
            # Step 5: Cleanup
            try:
                if temp_path.exists():
                    temp_path.unlink()
                    logger.info(f"STEP 5 COMPLETE: Cleaned up temporary file: {temp_path}")
                else:
                    logger.debug(f"STEP 5: Temporary file already removed: {temp_path}")
            except Exception as e:
                logger.warning(f"STEP 5 FAILED: Error cleaning up temporary file {temp_path}: {e}")
